%pip install pptx
%pip install prefect

from prefect import flow, task
from pptx import Presentation

@task
def generate_ppt(record_id: int):
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = f"Record {record_id}"
    slide.placeholders[1].text = "Generated via Quickbase + Prefect"
    
    file_path = f"record_{record_id}.pptx"
    ppt.save(file_path)
    return file_path

@flow(name="Generate PPT Flow")
def generate_ppt_flow(record_id: int):
    path = generate_ppt(record_id)
    return path
