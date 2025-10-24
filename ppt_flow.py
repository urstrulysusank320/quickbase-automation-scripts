from prefect_github import GitHub
from prefect import flow, task

# Load the GitHub repo block
github_block = GitHub.load("My GitHub Repo")

@task
def run_script():
    # Pull the file content from repo
    content = github_block.get_file_content("ppt_flow.py")  # path relative to repo root
    # Save it temporarily in the cloud environment
    with open("/tmp/ppt_flow.py", "w") as f:
        f.write(content)
    # Now import and run the script
    import ppt_flow
    ppt_flow.ppt_flow(title="QuickBase PPT")

@flow
def ppt_flow_trigger():
    run_script()

from pptx import Presentation

def ppt_flow(title="Default PPT"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    file_path = f"/tmp/{title}.pptx"
    prs.save(file_path)
    return file_path

